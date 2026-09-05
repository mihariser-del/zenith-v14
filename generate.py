import os
import csv
import io
import uuid
from datetime import datetime, timezone
from fastapi import APIRouter, Depends, Request, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy.ext.asyncio import AsyncSession
from database import get_db
from auth import get_current_user_from_cookie
from limits import check_limit

router = APIRouter(prefix="/api/generate", tags=["generate"])


class GenerateDocRequest(BaseModel):
    content: str
    filename: str = "document"
    format: str = "md"


@router.post("/document")
async def generate_document(req: GenerateDocRequest, request: Request, db: AsyncSession = Depends(get_db)):
    user = await get_current_user_from_cookie(request, db)
    # File editor/generator system: guests blocked entirely, logged-in (free)
    # users get 10 combined create/edit actions/day then a fixed 1h cooldown.
    await check_limit(user, db, "file_tool")
    content = req.content
    fmt = req.format.lower()
    filename = req.filename[:50]

    if fmt == "md":
        return StreamingResponse(
            io.BytesIO(content.encode("utf-8")),
            media_type="text/markdown",
            headers={"Content-Disposition": f'attachment; filename="{filename}.md"'}
        )

    elif fmt == "txt":
        return StreamingResponse(
            io.BytesIO(content.encode("utf-8")),
            media_type="text/plain",
            headers={"Content-Disposition": f'attachment; filename="{filename}.txt"'}
        )

    elif fmt == "csv":
        output = io.StringIO()
        writer = csv.writer(output)
        lines = content.split("\n")
        for line in lines:
            if line.strip():
                if "|" in line:
                    writer.writerow([cell.strip() for cell in line.split("|") if cell.strip()])
                else:
                    writer.writerow([line])
        output.seek(0)
        return StreamingResponse(
            io.BytesIO(output.getvalue().encode("utf-8")),
            media_type="text/csv",
            headers={"Content-Disposition": f'attachment; filename="{filename}.csv"'}
        )

    elif fmt == "html":
        html = f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8"><title>{filename}</title>
<style>
body {{ font-family: system-ui, sans-serif; max-width: 800px; margin: 40px auto; padding: 0 20px; line-height: 1.6; color: #333; }}
h1 {{ color: #0066ff; border-bottom: 2px solid #0066ff; padding-bottom: 10px; }}
pre {{ background: #f5f5f5; padding: 16px; border-radius: 8px; overflow-x: auto; }}
code {{ background: #f0f0f0; padding: 2px 6px; border-radius: 4px; font-size: 14px; }}
table {{ border-collapse: collapse; width: 100%; margin: 16px 0; }}
th, td {{ border: 1px solid #ddd; padding: 8px 12px; text-align: left; }}
th {{ background: #f5f5f5; }}
</style></head><body>
<pre>{content}</pre></body></html>"""
        return StreamingResponse(
            io.BytesIO(html.encode("utf-8")),
            media_type="text/html",
            headers={"Content-Disposition": f'attachment; filename="{filename}.html"'}
        )

    elif fmt == "pdf":
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet

            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=50, bottomMargin=50)
            styles = getSampleStyleSheet()
            elements = []

            for line in content.split("\n"):
                if line.startswith("# "):
                    elements.append(Paragraph(line[2:], styles['Title']))
                    elements.append(Spacer(1, 12))
                elif line.startswith("## "):
                    elements.append(Paragraph(line[3:], styles['Heading2']))
                    elements.append(Spacer(1, 8))
                elif line.startswith("### "):
                    elements.append(Paragraph(line[4:], styles['Heading3']))
                    elements.append(Spacer(1, 6))
                elif line.strip():
                    safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    elements.append(Paragraph(safe, styles['Normal']))
                    elements.append(Spacer(1, 4))

            doc.build(elements)
            buffer.seek(0)
            return StreamingResponse(
                buffer,
                media_type="application/pdf",
                headers={"Content-Disposition": f'attachment; filename="{filename}.pdf"'}
            )
        except ImportError:
            html = f"<html><body><pre>{content}</pre></body></html>"
            return StreamingResponse(
                io.BytesIO(html.encode("utf-8")),
                media_type="text/html",
                headers={"Content-Disposition": f'attachment; filename="{filename}.html"'}
            )

    elif fmt == "docx":
        try:
            from docx import Document
            from docx.shared import Pt
            doc = Document()
            for line in content.split("\n"):
                if line.startswith("# "):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith("## "):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith("### "):
                    doc.add_heading(line[4:], level=3)
                elif line.startswith("- "):
                    doc.add_paragraph(line[2:], style='List Bullet')
                elif line.strip():
                    doc.add_paragraph(line)

            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return StreamingResponse(
                buffer,
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": f'attachment; filename="{filename}.docx"'}
            )
        except ImportError:
            return StreamingResponse(
                io.BytesIO(content.encode("utf-8")),
                media_type="text/plain",
                headers={"Content-Disposition": f'attachment; filename="{filename}.txt"'}
            )

    elif fmt == "pptx":
        try:
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = filename
            title_frame.paragraphs[0].runs[0].font.size = Inches(0.3)
            title_frame.paragraphs[0].runs[0].font.bold = True
            # Content slides - split by headings
            slides_content = []
            current = []
            for line in content.split("\n"):
                if line.startswith("# ") and current:
                    slides_content.append("\n".join(current))
                    current = [line[2:]]
                else:
                    current.append(line)
            if current:
                slides_content.append("\n".join(current))
            for chunk in slides_content[:20]:  # max 20 slides
                if not chunk.strip():
                    continue
                s = prs.slides.add_slide(prs.slide_layouts[6])
                txBox = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                for line in chunk.split("\n"):
                    p = tf.add_paragraph()
                    p.text = line.replace("#", "").strip()
                    p.font.size = Inches(0.15)
                    if line.startswith("# "):
                        p.runs[0].font.bold = True
                        p.runs[0].font.size = Inches(0.2)
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return StreamingResponse(
                buffer,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f'attachment; filename="{filename}.pptx"'}
            )
        except ImportError:
            return StreamingResponse(
                io.BytesIO(content.encode("utf-8")),
                media_type="text/plain",
                headers={"Content-Disposition": f'attachment; filename="{filename}.txt"'}
            )
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PPTX generation failed: {str(e)}")

    elif fmt == "xlsx":
        try:
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = filename[:30]
            for r, line in enumerate(content.split("\n"), 1):
                if "|" in line:
                    cells = [c.strip() for c in line.split("|") if c.strip()]
                    for c, val in enumerate(cells, 1):
                        ws.cell(row=r, column=c, value=val)
                elif line.strip():
                    ws.cell(row=r, column=1, value=line)
            buffer = io.BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            return StreamingResponse(
                buffer,
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                headers={"Content-Disposition": f'attachment; filename="{filename}.xlsx"'}
            )
        except ImportError:
            return StreamingResponse(
                io.BytesIO(content.encode("utf-8")),
                media_type="text/plain",
                headers={"Content-Disposition": f'attachment; filename="{filename}.txt"'}
            )

    raise HTTPException(status_code=400, detail=f"Format not supported: {fmt}")
